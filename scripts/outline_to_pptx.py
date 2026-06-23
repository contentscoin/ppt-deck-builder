#!/usr/bin/env python3
"""Create a simple PPTX scaffold from slide_outline.json."""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def as_text(value: object) -> str:
    if isinstance(value, list):
        return "\n".join(f"- {item}" for item in value)
    if value is None:
        return ""
    return str(value)


def add_textbox(slide, x, y, w, h, text, size=18, bold=False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    return shape


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("outline")
    parser.add_argument("--out", required=True)
    args = parser.parse_args()

    outline = json.loads(Path(args.outline).read_text(encoding="utf-8"))
    slides = outline.get("slides", outline if isinstance(outline, list) else [])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for index, spec in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        role = spec.get("role", "claim")
        title = spec.get("title_claim") or spec.get("title") or f"Slide {index}"
        body = spec.get("supporting_points") or spec.get("body") or spec.get("notes") or ""
        source = spec.get("source_note") or spec.get("data_source") or ""

        if role == "cover":
            add_textbox(slide, 0.7, 1.7, 11.8, 1.4, title, size=44, bold=True)
            add_textbox(slide, 0.75, 3.25, 8.5, 1.0, as_text(body), size=20)
        elif role == "section_divider":
            add_textbox(slide, 0.7, 2.45, 11.8, 1.0, title, size=38, bold=True)
            add_textbox(slide, 0.75, 3.55, 9.5, 0.8, as_text(body), size=18)
        else:
            add_textbox(slide, 0.55, 0.45, 12.1, 0.75, title, size=30, bold=True)
            add_textbox(slide, 0.75, 1.55, 7.1, 4.6, as_text(body), size=18)
            visual = spec.get("primary_visual", "")
            if visual:
                add_textbox(slide, 8.25, 1.55, 4.3, 3.8, f"Visual: {as_text(visual)}", size=14)

        if source:
            add_textbox(slide, 0.55, 7.05, 12.2, 0.25, f"Source: {as_text(source)}", size=7)

    prs.save(args.out)
    print(json.dumps({"ok": True, "pptx": args.out, "slides": len(slides)}, indent=2))


if __name__ == "__main__":
    main()
